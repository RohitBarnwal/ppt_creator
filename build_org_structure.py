import os
import sys
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR

def build_org_presentation(excel_path, output_pptx_path):
    print(f"Reading Excel: {excel_path}")
    if not os.path.exists(excel_path):
        print(f"Error: Excel file '{excel_path}' does not exist!")
        return False
        
    xl = pd.ExcelFile(excel_path)
    
    # 1. Robust Sheet Detection & Concatenation
    dfs = []
    for sheet in xl.sheet_names:
        df_temp = xl.parse(sheet)
        cols_lower = [str(c).lower().strip() for c in df_temp.columns]
        if any('employee name' in c or 'employee_name' in c for c in cols_lower) and \
           any('manager name' in c or 'manager_name' in c for c in cols_lower):
            print(f"Discovered employee sheet: '{sheet}' with {df_temp.shape[0]} rows")
            df_temp.columns = [str(c).strip() for c in df_temp.columns]
            dfs.append(df_temp)
            
    if not dfs:
        print("Error: Could not find any worksheets with employee databases!")
        return False
        
    # Concatenate all matching tables to construct a complete, unified database
    df = pd.concat(dfs, ignore_index=True, sort=False)
    print(f"Unified database shape: {df.shape}")
    
    # 2. Extract and Normalize Column Names
    name_col = [c for c in df.columns if 'employee name' in c.lower() or c.lower() == 'name'][0]
    manager_col = [c for c in df.columns if 'l1 manager name' in c.lower() or 'manager name' in c.lower()][0]
    designation_col = [c for c in df.columns if 'designation' in c.lower()][0]
    
    # Prioritize sub-department over general BU/SBU
    sub_dept_col = None
    for kw in ['sub department', 'sub_department', 'department', 'sbu', 'bu']:
        matches = [c for c in df.columns if kw in c.lower()]
        if matches:
            sub_dept_col = matches[0]
            break
    if not sub_dept_col:
        sub_dept_col = df.columns[-1]
    
    df[name_col] = df[name_col].astype(str).str.strip()
    df[manager_col] = df[manager_col].astype(str).str.strip()
    df[designation_col] = df[designation_col].astype(str).str.strip()
    df[sub_dept_col] = df[sub_dept_col].astype(str).str.strip()
    
    # 3. Recursive Team Headcount Traversal
    def get_recursive_hc(person_name, visited=None):
        if visited is None:
            visited = set()
        if person_name in visited:
            return 0
        visited.add(person_name)
        
        # Find direct reports
        directs = df[df[manager_col] == person_name]
        count = len(directs)
        for idx, row in directs.iterrows():
            rname = row[name_col]
            if rname != person_name:
                count += get_recursive_hc(rname, visited)
        return count

    # 4. Auto-detect Root Leader
    candidate_managers = [m for m in df[manager_col].unique() if m and str(m).lower() not in ['nan', 'none', 'unknown']]
    manager_team_sizes = {mgr: get_recursive_hc(mgr) for mgr in candidate_managers}
    root_manager = max(manager_team_sizes, key=manager_team_sizes.get)
    print(f"Auto-detected root manager: '{root_manager}' with a total span of {manager_team_sizes[root_manager]} employees.")
    
    # 5. Determine Top Senior Chain
    is_control_functions = "sharma" in root_manager.lower()
    
    senior_chain = []
    if is_control_functions:
        # For Control Functions, Vijay Shekhar Sharma is the sole top lead
        senior_chain = [
            {"name": "Vijay Shekhar Sharma", "title": "MD & CEO"}
        ]
    else:
        # For Gagan's span, trace upper managers
        root_row = df[df[name_col] == root_manager]
        if not root_row.empty:
            row = root_row.iloc[0]
            for idx in range(1, 7):
                m_col = f"L{idx} Manager Name"
                if m_col in df.columns and pd.notna(row[m_col]) and str(row[m_col]).strip().lower() not in ['nan', 'none', '']:
                    m_name = str(row[m_col]).strip()
                    m_des = "MD & CEO" if "sharma" in m_name.lower() else ""
                    if not m_des:
                        m_row = df[df[name_col] == m_name]
                        m_des = m_row.iloc[0][designation_col] if not m_row.empty else "Executive"
                    senior_chain.append({"name": m_name, "title": m_des})
        
        if not senior_chain:
            senior_chain = [
                {"name": "Vijay Shekhar Sharma", "title": "MD & CEO"},
                {"name": "Deependra Singh Rathore", "title": "COO - Tech, Product and Ops"}
            ]
        else:
            senior_chain = senior_chain[::-1]
        
    print("Detected Senior Chain:")
    for s in senior_chain:
        print(f"  - {s['name']} ({s['title']})")
        
    # Get direct reports under the root manager (the HODs)
    hods = df[df[manager_col] == root_manager]
    print(f"Found {len(hods)} HODs reporting directly to {root_manager}:")
    for _, h in hods.iterrows():
        print(f"  - {h[name_col]} | Designation: {h[designation_col]}")
        
    # 6. Initialize Presentation
    prs = Presentation()
    prs.slide_width = Inches(10.0)
    prs.slide_height = Inches(5.62)
    blank_layout = prs.slide_layouts[6]
    
    # Exact Color Palette Theme
    color_vss = RGBColor(7, 55, 99)       # Dark Navy Blue (#073763)
    color_root = RGBColor(11, 83, 148)    # Medium Dark Blue (#0B5394)
    color_hod = RGBColor(61, 133, 198)    # Medium Blue (#3D85C6)
    color_manager = RGBColor(95, 157, 209) # Light Blue (#5F9DD1)
    color_report = RGBColor(143, 188, 224) # Very Light Blue (#8FC0E0)
    
    # Header Title Mapping (to match reference titles perfectly)
    header_map = {
        "compliance": "PPSL – AML & Compliance",
        "information and cyber security": "PPSL – Infosec",
        "internal audit": "PPSL – Internal Audit",
        "legal litigation": "PPSL – Legal",
        "risk management": "PPSL – Ops & Fraud Risk"
    }
    
    # Helper to add cards
    def add_employee_box(slide, name, title, x, y, cx, cy, fill_color, text_color=RGBColor(255, 255, 255), is_bold=False):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_top = Inches(0.01)
        tf.margin_bottom = Inches(0.01)
        tf.margin_left = Inches(0.01)
        tf.margin_right = Inches(0.01)
        
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = 'Arial'
        p.font.size = Pt(7.0)
        p.font.bold = is_bold
        p.font.color.rgb = text_color
        p.alignment = 1 # Center
        
        if title:
            p2 = tf.add_paragraph()
            p2.text = title
            p2.font.name = 'Arial'
            p2.font.size = Pt(6.0)
            p2.font.color.rgb = text_color
            p2.alignment = 1 # Center
        return shape

    # Helper to draw a perfectly straight vertical or horizontal line (Bold, zero turns, zero slants)
    def add_straight_line(slide, x1, y1, x2, y2):
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, int(x1), int(y1), int(x2), int(y2))
        connector.line.color.rgb = RGBColor(127, 127, 127)
        connector.line.width = Pt(1.5) # BOLD lines!
        return connector

    # Helper to draw a perfect top-down fork (Stem -> Bar -> Drops) with absolute zero slants and clean 90-degree drops
    def add_perfect_fork(slide, parent_shape, child_shapes):
        if not child_shapes:
            return
            
        parent_x = parent_shape.left + parent_shape.width / 2
        parent_bottom = parent_shape.top + parent_shape.height
        child_top = child_shapes[0].top
        
        # Draw a single direct vertical line if only 1 child
        if len(child_shapes) == 1:
            child_x = child_shapes[0].left + child_shapes[0].width / 2
            # Force identical x-centers to guarantee absolutely vertical alignment and zero slants
            add_straight_line(slide, parent_x, parent_bottom, parent_x, child_top)
            return
            
        # Draw perfect 3-segment fork structure: Vertical Stem -> Horizontal Bar -> Vertical Drops
        y_bar = parent_bottom + (child_top - parent_bottom) / 2
        
        # 1. Stem
        add_straight_line(slide, parent_x, parent_bottom, parent_x, y_bar)
        
        # 2. Horizontal Bar
        x_coords = [s.left + s.width / 2 for s in child_shapes]
        x_min = min(x_coords)
        x_max = max(x_coords)
        add_straight_line(slide, x_min, y_bar, x_max, y_bar)
        
        # 3. Drops
        for s in child_shapes:
            cx = s.left + s.width / 2
            add_straight_line(slide, cx, y_bar, cx, child_top)

    # ----------------------------------------------------
    # SLIDE 1: Cover Slide
    # ----------------------------------------------------
    slide_t = prs.slides.add_slide(blank_layout)
    bg_shape = slide_t.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(10.0), Inches(5.62))
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = color_vss
    bg_shape.line.fill.background()
    
    tf_title = bg_shape.text_frame
    p_main = tf_title.paragraphs[0]
    p_main.text = "PPSL - Control Functions Org Structure" if is_control_functions else "Paytm Payments Services Limited"
    p_main.font.name = 'Calibri'
    p_main.font.size = Pt(28)
    p_main.font.bold = True
    p_main.alignment = 1
    
    if not is_control_functions:
        p_sub = tf_title.add_paragraph()
        p_sub.text = f"\n{root_manager}'s Span Org Structure\nJune 2026"
        p_sub.font.name = 'Calibri'
        p_sub.font.size = Pt(20)
        p_sub.alignment = 1

    # ----------------------------------------------------
    # SLIDE 2: Primary HOD Overview Grid (Only for Gagan's Span)
    # ----------------------------------------------------
    if not is_control_functions:
        slide_o = prs.slides.add_slide(blank_layout)
        
        # Slide Title
        title_box = slide_o.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.2), Inches(0.35))
        title_box.fill.solid()
        title_box.fill.fore_color.rgb = color_vss
        title_box.line.fill.background()
        p_t = title_box.text_frame.paragraphs[0]
        p_t.text = f"{root_manager}’s HODs Overview"
        p_t.font.name = 'Calibri'
        p_t.font.size = Pt(13)
        p_t.font.bold = True
        p_t.font.color.rgb = RGBColor(255, 255, 255)
        
        # Upper Hierarchy Cards
        vss_card = add_employee_box(slide_o, senior_chain[0]['name'], senior_chain[0]['title'], Inches(4.1), Inches(0.15), Inches(1.8), Inches(0.38), color_vss, is_bold=True)
        coo_card = add_employee_box(slide_o, senior_chain[1]['name'], senior_chain[1]['title'], Inches(4.1), Inches(0.60), Inches(1.8), Inches(0.38), color_vss, is_bold=True)
        root_card = add_employee_box(slide_o, root_manager, f"Root | HC - {manager_team_sizes[root_manager]}", Inches(4.1), Inches(1.05), Inches(1.8), Inches(0.42), color_root, is_bold=True)
        add_perfect_fork(slide_o, vss_card, [coo_card])
        add_perfect_fork(slide_o, coo_card, [root_card])
        
        hod_list = sorted(list(hods.iterrows()), key=lambda x: get_recursive_hc(x[1][name_col]), reverse=True)
        num_hods = len(hod_list)
        
        row0_count = (num_hods + 1) // 2
        row1_count = num_hods - row0_count
        row0_span = Inches(9.6)
        col_width, col_height = Inches(1.3), Inches(0.55)
        
        # Row 0
        gap0 = row0_span / row0_count if row0_count > 0 else Inches(1)
        for j in range(row0_count):
            _, h_row = hod_list[j]
            hname, htitle, hsub = h_row[name_col], h_row[designation_col], h_row[sub_dept_col]
            hhc = get_recursive_hc(hname)
            col_center = Inches(0.2) + (j + 0.5) * gap0
            hx = col_center - col_width / 2
            add_employee_box(slide_o, hname, f"{htitle}\n{hsub}\nHC - {hhc}", hx, Inches(1.8), col_width, col_height, color_hod, is_bold=True)
            
        # Row 1
        if row1_count > 0:
            gap1 = row0_span / row1_count
            for j in range(row1_count):
                _, h_row = hod_list[row0_count + j]
                hname, htitle, hsub = h_row[name_col], h_row[designation_col], h_row[sub_dept_col]
                hhc = get_recursive_hc(hname)
                col_center = Inches(0.2) + (j + 0.5) * gap1
                hx = col_center - col_width / 2
                add_employee_box(slide_o, hname, f"{htitle}\n{hsub}\nHC - {hhc}", hx, Inches(3.3), col_width, col_height, color_hod, is_bold=True)
                
        hc_box_o = slide_o.shapes.add_textbox(Inches(7.8), Inches(5.1), Inches(2.0), Inches(0.4))
        hc_box_o.text_frame.paragraphs[0].text = f"Total HC: {manager_team_sizes[root_manager]}\nDR- Direct Reportee"
        hc_box_o.text_frame.paragraphs[0].font.name = 'Arial'
        hc_box_o.text_frame.paragraphs[0].font.size = Pt(8)
        hc_box_o.text_frame.paragraphs[0].font.bold = True

    # ----------------------------------------------------
    # SLIDES 3+: Dedicated HOD team slide (goes up to Level 2 only)
    # ----------------------------------------------------
    for idx, h_row in hods.iterrows():
        hname, htitle, hsub = h_row[name_col], h_row[designation_col], h_row[sub_dept_col]
        team_size = get_recursive_hc(hname)
        
        if team_size == 0:
            continue
            
        directs = df[df[manager_col] == hname]
        num_directs = len(directs)
        
        if num_directs == 0:
            continue
            
        # Split direct reports into chunks of at most 6 reports per slide
        chunk_size = 6
        directs_list = list(directs.iterrows())
        chunks = [directs_list[i:i + chunk_size] for i in range(0, len(directs_list), chunk_size)]
        
        for chunk_idx, chunk in enumerate(chunks):
            slide_h = prs.slides.add_slide(blank_layout)
            
            # Format Header Slide Title
            header_suffix = " (Contd.)" if chunk_idx > 0 else ""
            formatted_header = header_map.get(hsub.lower().strip(), f"PPSL – {hsub}" if is_control_functions else f"PG Tech – {hsub}") + header_suffix
            
            title_box_h = slide_h.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(3.2), Inches(0.35))
            title_box_h.fill.solid()
            title_box_h.fill.fore_color.rgb = color_vss
            title_box_h.line.fill.background()
            p_th = title_box_h.text_frame.paragraphs[0]
            p_th.text = formatted_header
            p_th.font.name = 'Calibri'
            p_th.font.size = Pt(13)
            p_th.font.bold = True
            p_th.font.color.rgb = RGBColor(255, 255, 255)
            
            # Senior chain placement
            if is_control_functions:
                # ONLY VSS on top for Control Functions
                v_card = add_employee_box(slide_h, senior_chain[0]['name'], senior_chain[0]['title'], Inches(4.15), Inches(0.20), Inches(1.7), Inches(0.40), color_vss, is_bold=True)
                h_card = add_employee_box(slide_h, hname, f"{htitle}\nDR-{get_recursive_hc(hname)}", Inches(4.15), Inches(0.85), Inches(1.7), Inches(0.45), color_hod, is_bold=True)
                add_perfect_fork(slide_h, v_card, [h_card])
            else:
                v_card = add_employee_box(slide_h, senior_chain[0]['name'], senior_chain[0]['title'], Inches(4.1), Inches(0.15), Inches(1.8), Inches(0.38), color_vss, is_bold=True)
                c_card = add_employee_box(slide_h, senior_chain[1]['name'], senior_chain[1]['title'], Inches(4.1), Inches(0.60), Inches(1.8), Inches(0.38), color_vss, is_bold=True)
                h_card = add_employee_box(slide_h, hname, f"{htitle}\nHC - {team_size}", Inches(4.1), Inches(1.05), Inches(1.8), Inches(0.42), color_hod, is_bold=True)
                add_perfect_fork(slide_h, v_card, [c_card])
                add_perfect_fork(slide_h, c_card, [h_card])
                
            # Lay out the direct reports in this chunk horizontally
            row_items = len(chunk)
            col_width_h = Inches(1.4) # Always 1.4" since we have at most 6 items per slide
            col_height_h = Inches(0.50)
            
            y_pos = Inches(1.8) if not is_control_functions else Inches(1.7)
            row_span = Inches(9.6)
            gap = row_span / row_items
            
            report_cards = []
            for j in range(row_items):
                _, d_row = chunk[j]
                dname, dtitle, dsub = d_row[name_col], d_row[designation_col], d_row[sub_dept_col]
                dhc = get_recursive_hc(dname)
                col_center = Inches(0.2) + (j + 0.5) * gap
                dx = col_center - col_width_h / 2
                
                if dhc > 0:
                    # Manager rollup: Name, Title, and DR-count
                    card_txt = f"{dtitle}\nDR- {dhc}"
                    card = add_employee_box(slide_h, dname, card_txt, dx, y_pos, col_width_h, col_height_h, color_manager, is_bold=True)
                    
                    # Generate vertical sub-stack cards for sub-departments
                    sub_reports = df[df[manager_col] == dname]
                    if not sub_reports.empty:
                        # Group by sub-department
                        group_counts = sub_reports[sub_dept_col].value_counts()
                        groups = [f"{dept} - {count}" for dept, count in group_counts.items()]
                        
                        # Group into pairs (at most 2 groups per sub-card)
                        card_texts = []
                        for k in range(0, len(groups), 2):
                            if k + 1 < len(groups):
                                card_texts.append(f"{groups[k]}\n{groups[k+1]}")
                            else:
                                card_texts.append(groups[k])
                                
                        if card_texts:
                            # 1. Stem of the left-aligned fork
                            x_stem = dx + Inches(0.12)
                            y_stem_start = y_pos + col_height_h
                            y_stem_end = (y_pos + col_height_h + Inches(0.22)) + (len(card_texts) - 1) * Inches(0.72) + col_height_h / 2
                            add_straight_line(slide_h, x_stem, y_stem_start, x_stem, y_stem_end)
                            
                            for i_card, card_text in enumerate(card_texts):
                                sub_y = (y_pos + col_height_h + Inches(0.22)) + i_card * Inches(0.72)
                                sub_x = dx + Inches(0.25)
                                col_width_sub = Inches(1.15)
                                
                                # Add sub-card shape (text centered vertically by leaving title empty)
                                add_employee_box(
                                    slide_h, 
                                    card_text, 
                                    "", 
                                    sub_x, 
                                    sub_y, 
                                    col_width_sub, 
                                    col_height_h, 
                                    color_manager
                                )
                                
                                # 2. Horizontal branch connecting stem to sub-card
                                y_branch = sub_y + col_height_h / 2
                                add_straight_line(slide_h, x_stem, y_branch, sub_x, y_branch)
                else:
                    # Individual contributor card
                    card = add_employee_box(slide_h, dname, dtitle, dx, y_pos, col_width_h, col_height_h, color_report)
                    
                report_cards.append(card)
                
            # Connect the parent HOD to the entire row of direct reports using our perfect zero-slant fork structure
            add_perfect_fork(slide_h, h_card, report_cards)
                
            # Total headcount tally matching reference
            hc_box_h = slide_h.shapes.add_textbox(Inches(7.8), Inches(5.1), Inches(2.0), Inches(0.4))
            tf_hc_h = hc_box_h.text_frame
            p_hc_h = tf_hc_h.paragraphs[0]
            
            # Use exact text format (Active + SNP) for Control Functions
            status_suffix = "(Active + SNP)" if is_control_functions else "(Active)"
            p_hc_h.text = f"Total HC- {team_size} {status_suffix}\nDR- Direct Reportee" if is_control_functions else f"HC : {team_size} (Active + Intern)\nDR- Direct Reportee"
            p_hc_h.font.name = 'Arial'
            p_hc_h.font.size = Pt(8)
            p_hc_h.font.bold = True
        
    print(f"Saving final Presentation to: {output_pptx_path}")
    prs.save(output_pptx_path)
    print("Org structure generation complete!")
    return True

if __name__ == "__main__":
    excel_file = "_List Of Employees sheet __ Gagan' span __ June'26.xlsx"
    output_pptx = "Automated_PPSL_Org_Structure_June26.pptx"
    build_org_presentation(excel_file, output_pptx)
